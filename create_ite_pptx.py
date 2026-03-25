"""
Create ITE topic PowerPoints.
- For each question: Question + options -> Poll responses (bar chart) -> Critique
- Results slide has an editable chart (A–E); enter vote counts after your poll.
- Grouped by topic and subtopic. AFP-style: Montserrat, Segoe UI, grey box, 16:9
"""

import json
import re
from pathlib import Path
from collections import defaultdict

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

SCRIPT_DIR = Path(__file__).parent
DATA_FILE = SCRIPT_DIR / "ite_data.json"
OUTPUT_DIR = SCRIPT_DIR / "ITE_PowerPoints"

# AFP-style layout (from create_ppt_from_web.py)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.5)
TOP_MARGIN = Inches(0.42)
FONT_TITLE = "Montserrat"
FONT_BODY = "Segoe UI"
GREY_BOX_BG = RGBColor(0xEA, 0xEA, 0xEA) if HAS_PPTX else None
GREY_BOX_TEXT = RGBColor(0x76, 0x76, 0x76) if HAS_PPTX else None
BLURB_COLOR = RGBColor(0x55, 0x55, 0x55) if HAS_PPTX else None
GREY_BOX_HEIGHT = Inches(0.2)
TITLE_SIZE = Pt(28)
SUBTOPIC_SIZE = Pt(14)  # Small heading for sub-topic
BODY_SIZE = Pt(16)
QUESTION_STEM_SIZE = Pt(18)
OPTION_SIZE = Pt(14)
CRITIQUE_SIZE = Pt(14)


def _sanitize(s: str, max_len: int = 2000) -> str:
    """Remove chars that can corrupt PPTX."""
    if not s:
        return ""
    s = "".join(c for c in str(s) if ord(c) >= 32 or c in "\n\r\t")[:max_len]
    return s


def _wrap_text(text: str, max_chars: int = 80) -> str:
    """Soft wrap long lines for readability."""
    words = text.split()
    lines = []
    current = []
    current_len = 0
    for w in words:
        if current_len + len(w) + 1 <= max_chars:
            current.append(w)
            current_len += len(w) + 1
        else:
            if current:
                lines.append(" ".join(current))
            current = [w]
            current_len = len(w) + 1
    if current:
        lines.append(" ".join(current))
    return "\n".join(lines)


def _add_grey_box(slide, label: str, y_in: float, box_w) -> float:
    """Add grey box header. Returns y after box."""
    try:
        lbl_box = slide.shapes.add_textbox(MARGIN, Inches(y_in), box_w, GREY_BOX_HEIGHT)
        tf = lbl_box.text_frame
        tf.word_wrap = False
        tf.text = label.upper()
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = GREY_BOX_TEXT
        p.font.name = FONT_TITLE
    except Exception:
        pass
    return y_in + 0.2 + 0.08


def _add_poll_question_slide(prs, blank_layout, q: dict, topic: str, subtopic: str, year: int) -> None:
    """Question slide: stem + all answer choices."""
    slide = prs.slides.add_slide(blank_layout)
    box_w = SLIDE_WIDTH - 2 * MARGIN
    y = 0.42
    
    # Grey box: TOPIC | SUBTOPIC
    y = _add_grey_box(slide, f"{topic} | {subtopic}", y, box_w)
    
    # Small sub-topic heading
    sub_box = slide.shapes.add_textbox(MARGIN, Inches(y), box_w, Inches(0.35))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = f"Question {q['number']} ({year})"
    p_sub.font.size = SUBTOPIC_SIZE
    p_sub.font.name = FONT_TITLE
    try:
        p_sub.font.color.rgb = GREY_BOX_TEXT
    except Exception:
        pass
    y += 0.4
    
    # Question stem
    stem = _sanitize(q["stem"], 1200)
    stem_wrapped = _wrap_text(stem, 90)
    stem_box = slide.shapes.add_textbox(MARGIN, Inches(y), box_w, Inches(2.2))
    tf_stem = stem_box.text_frame
    tf_stem.word_wrap = True
    p_stem = tf_stem.paragraphs[0]
    p_stem.text = stem_wrapped
    p_stem.font.size = QUESTION_STEM_SIZE
    p_stem.font.name = FONT_BODY
    y += 2.4
    
    # Options A-E
    for letter in ["A", "B", "C", "D", "E"]:
        if letter in q.get("options", {}):
            opt_text = _sanitize(q["options"][letter], 300)
            opt_box = slide.shapes.add_textbox(MARGIN, Inches(y), box_w, Inches(0.5))
            tf_opt = opt_box.text_frame
            tf_opt.word_wrap = True
            p_opt = tf_opt.paragraphs[0]
            p_opt.text = f"{letter}) {opt_text}"
            p_opt.font.size = OPTION_SIZE
            p_opt.font.name = FONT_BODY
            y += 0.5


def _option_letters(q: dict) -> list[str]:
    opts = q.get("options") or {}
    return [L for L in ("A", "B", "C", "D", "E") if L in opts]


def _add_poll_responses_slide(prs, blank_layout, q: dict, topic: str, subtopic: str, year: int) -> None:
    """Horizontal bar chart for poll vote counts; correct choice bar in green when known."""
    slide = prs.slides.add_slide(blank_layout)
    box_w = SLIDE_WIDTH - 2 * MARGIN
    y = 0.42

    y = _add_grey_box(slide, f"{topic} | {subtopic}", y, box_w)

    title_box = slide.shapes.add_textbox(MARGIN, Inches(y), box_w, Inches(0.55))
    tf_t = title_box.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = f"Poll responses — Question {q['number']} ({year})"
    p_t.font.size = TITLE_SIZE
    p_t.font.name = FONT_TITLE
    p_t.font.bold = True
    y += 0.62

    hint = slide.shapes.add_textbox(MARGIN, Inches(y), box_w, Inches(0.38))
    hp = hint.text_frame.paragraphs[0]
    hp.text = (
        "Automate: fill poll_results.csv (see export_poll_template.py), then run "
        "apply_poll_results.py --csv poll_results.csv --pptx this deck"
    )
    hp.font.size = Pt(11)
    hp.font.name = FONT_BODY
    try:
        hp.font.color.rgb = BLURB_COLOR
    except Exception:
        pass
    y += 0.42

    letters_ord = _option_letters(q)
    if not letters_ord:
        return
    zeros = tuple(0 for _ in letters_ord)

    chart_data = CategoryChartData()
    chart_data.categories = letters_ord
    chart_data.add_series("Responses", zeros)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        MARGIN,
        Inches(y),
        box_w,
        Inches(4.6),
        chart_data,
    )
    chart = graphic_frame.chart
    try:
        chart.value_axis.has_major_gridlines = False
    except Exception:
        pass
    try:
        chart.has_legend = False
    except Exception:
        pass

    correct = (q.get("answer") or "").strip().upper()[:1]
    if correct in letters_ord:
        try:
            series = chart.plots[0].series[0]
            idx = letters_ord.index(correct)
            pt = series.points[idx]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = RGBColor(0x28, 0xA7, 0x45)
        except Exception:
            pass


def _add_critique_slide(prs, blank_layout, q: dict, topic: str, subtopic: str, year: int) -> None:
    """Critique slide - answer and explanation."""
    slide = prs.slides.add_slide(blank_layout)
    box_w = SLIDE_WIDTH - 2 * MARGIN
    y = 0.42
    
    y = _add_grey_box(slide, f"{topic} | {subtopic}", y, box_w)
    
    # Header: Answer
    ans = q.get("answer", "?")
    header_box = slide.shapes.add_textbox(MARGIN, Inches(y), box_w, Inches(0.6))
    tf_h = header_box.text_frame
    tf_h.word_wrap = True
    p_h = tf_h.paragraphs[0]
    p_h.text = f"Answer: {ans} — Question {q['number']} ({year})"
    p_h.font.size = TITLE_SIZE
    p_h.font.name = FONT_TITLE
    p_h.font.bold = True
    y += 0.7
    
    # Critique text
    critique = q.get("critique") or "No critique available."
    critique = _sanitize(critique, 2500)
    critique_box = slide.shapes.add_textbox(MARGIN, Inches(y), box_w, Inches(5.8))
    tf_c = critique_box.text_frame
    tf_c.word_wrap = True
    tf_c.auto_size = None
    p_c = tf_c.paragraphs[0]
    p_c.text = critique
    p_c.font.size = CRITIQUE_SIZE
    p_c.font.name = FONT_BODY
    try:
        p_c.font.color.rgb = BLURB_COLOR
    except Exception:
        pass


def _add_intro_slide(prs, blank_layout, topic_name: str) -> None:
    """Title slide for the topic."""
    slide = prs.slides.add_slide(blank_layout)
    box_w = SLIDE_WIDTH - 2 * MARGIN
    title_y = (SLIDE_HEIGHT.inches / 2) - 0.8
    
    title_box = slide.shapes.add_textbox(MARGIN, Inches(title_y), box_w, Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic_name
    p.font.size = Pt(44)
    p.font.name = FONT_TITLE
    p.font.bold = True
    if PP_ALIGN:
        p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(MARGIN, Inches(title_y + 1.0), box_w, Inches(0.5))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "ITE Questions by Topic — Chiefs 2026-2027"
    p2.font.size = Pt(16)
    p2.font.name = FONT_BODY
    if PP_ALIGN:
        p2.alignment = PP_ALIGN.CENTER
    try:
        p2.font.color.rgb = BLURB_COLOR
    except Exception:
        pass


def create_topic_pptx(questions: list[dict], topic: str, output_path: Path) -> None:
    """Create one PowerPoint for a topic, grouped by subtopic."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]
    
    _add_intro_slide(prs, blank_layout, topic)

    # Group questions by subtopic (for ordering)
    by_subtopic: dict[str, list[dict]] = defaultdict(list)
    for q in questions:
        sub = q.get("subtopic") or "Other"
        by_subtopic[sub].append(q)
    
    # Sort subtopics by first question number
    def first_num(items):
        return min(item["number"] for item in items) if items else 0
    
    subtopic_order = sorted(by_subtopic.keys(), key=lambda s: first_num(by_subtopic[s]))
    
    for subtopic in subtopic_order:
        items = sorted(by_subtopic[subtopic], key=lambda q: (q["year"], q["number"]))
        for q in items:
            _add_poll_question_slide(prs, blank_layout, q, topic, subtopic, q["year"])
            _add_poll_responses_slide(prs, blank_layout, q, topic, subtopic, q["year"])
            _add_critique_slide(prs, blank_layout, q, topic, subtopic, q["year"])
    
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main():
    if not HAS_PPTX:
        print("Install python-pptx: pip install python-pptx")
        return 1
    
    data_path = DATA_FILE
    if not data_path.exists():
        data_path = SCRIPT_DIR / "ite_data_sample.json"
    if not data_path.exists():
        print(f"Run extract_ite.py first to create {DATA_FILE}")
        return 1
    
    with open(data_path, encoding="utf-8") as f:
        all_questions = json.load(f)
    
    # Group by topic
    by_topic: dict[str, list[dict]] = defaultdict(list)
    for q in all_questions:
        topic = q.get("topic") or "General Medicine"
        by_topic[topic].append(q)
    
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    for topic in sorted(by_topic.keys()):
        questions = by_topic[topic]
        safe_name = re.sub(r'[^\w\s-]', '', topic).strip().replace(" ", "_")
        out_path = OUTPUT_DIR / f"{safe_name}.pptx"
        create_topic_pptx(questions, topic, out_path)
        print(f"Created: {out_path} ({len(questions)} questions)")
    
    print(f"\nDone. Output in {OUTPUT_DIR}")
    return 0


if __name__ == "__main__":
    exit(main())
