"""
Apply poll vote counts to every "Poll responses" bar chart in ITE topic decks.

Reads a CSV (or JSON) with counts per answer letter, matches slides by
"Poll responses — Question N (Year)", updates charts and re-colors the correct answer.

CSV format (header required):
  year,number,A,B,C,D,E

JSON format:
  {"2024": {"1": {"A": 5, "B": 3, "C": 2, "D": 0, "E": 1}}}
  or list: [{"year": 2024, "number": 1, "votes": {"A": 5, ...}}, ...]

Usage:
  python apply_poll_results.py --csv poll_results.csv
  python apply_poll_results.py --csv poll_results.csv --pptx ITE_PowerPoints/Cardiology.pptx
  python apply_poll_results.py --csv poll_results.csv --all   # every *.pptx in ITE_PowerPoints
"""

from __future__ import annotations

import argparse
import csv
import json
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

SCRIPT_DIR = Path(__file__).parent
DATA_JSON = SCRIPT_DIR / "ite_data.json"
PPTS_DIR = SCRIPT_DIR / "ITE_PowerPoints"

# Title may use em dash, hyphen, or odd encoding from PowerPoint
TITLE_RE = re.compile(
    r"Poll responses\s*.*?\s*Question\s+(\d+)\s*\((\d{4})\)",
    re.IGNORECASE | re.DOTALL,
)


def load_answer_key() -> dict[tuple[int, int], str]:
    """(year, number) -> correct letter."""
    if not DATA_JSON.exists():
        return {}
    rows = json.loads(DATA_JSON.read_text(encoding="utf-8"))
    out = {}
    for q in rows:
        y, n = q.get("year"), q.get("number")
        a = (q.get("answer") or "").strip().upper()[:1]
        if y is not None and n is not None and a in "ABCDE":
            out[(int(y), int(n))] = a
    return out


def load_counts_csv(path: Path) -> dict[tuple[int, int], dict[str, int]]:
    with path.open(newline="", encoding="utf-8-sig") as f:
        r = csv.DictReader(f)
        if not r.fieldnames:
            raise ValueError("Empty CSV")
        colmap = {c.strip().lower(): c for c in r.fieldnames}
        def cell(row, name: str) -> str:
            k = colmap.get(name.lower())
            return (row.get(k) or "").strip() if k else ""
        out: dict[tuple[int, int], dict[str, int]] = {}
        for row in r:
            y = int(float(cell(row, "year")))
            n = int(float(cell(row, "number")))
            votes: dict[str, int] = {}
            for L in "ABCDE":
                raw = cell(row, L)
                votes[L] = int(float(raw)) if raw != "" else 0
            out[(y, n)] = votes
    return out


def load_counts_json(path: Path) -> dict[tuple[int, int], dict[str, int]]:
    raw = json.loads(path.read_text(encoding="utf-8"))
    if isinstance(raw, dict) and "polls" in raw:
        raw = raw["polls"]
    out: dict[tuple[int, int], dict[str, int]] = {}
    if isinstance(raw, list):
        for item in raw:
            y = int(item["year"])
            n = int(item["number"])
            v = item.get("votes") or item
            votes = {L: int(v.get(L, v.get(L.lower(), 0))) for L in "ABCDE"}
            out[(y, n)] = votes
        return out
    # nested dict year -> num -> votes
    for ys, nums in raw.items():
        y = int(ys)
        for ns, vdict in nums.items():
            n = int(ns)
            votes = {L: int(vdict.get(L, vdict.get(L.lower(), 0))) for L in "ABCDE"}
            out[(y, n)] = votes
    return out


def load_counts(path: Path) -> dict[tuple[int, int], dict[str, int]]:
    if path.suffix.lower() == ".json":
        return load_counts_json(path)
    return load_counts_csv(path)


def _slide_title_blob(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            t = (para.text or "").strip()
            if t:
                parts.append(t)
    return "\n".join(parts)


def _apply_chart(
    chart,
    categories: list[str],
    votes: dict[str, int],
    correct: str | None,
) -> None:
    vals = tuple(votes.get(c, 0) for c in categories)
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("Responses", vals)
    chart.replace_data(cd)
    if correct and correct in categories:
        try:
            idx = categories.index(correct)
            pt = chart.plots[0].series[0].points[idx]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = RGBColor(0x28, 0xA7, 0x45)
        except Exception:
            pass


def process_presentation(
    pptx_path: Path,
    counts: dict[tuple[int, int], dict[str, int]],
    answer_key: dict[tuple[int, int], str],
) -> tuple[int, int]:
    prs = Presentation(str(pptx_path))
    updated = 0
    missing = 0
    for slide in prs.slides:
        blob = _slide_title_blob(slide)
        m = TITLE_RE.search(blob)
        if not m:
            continue
        qnum, year = int(m.group(1)), int(m.group(2))
        key = (year, qnum)
        if key not in counts:
            missing += 1
            continue
        chart_shape = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_shape = shape
                break
        if chart_shape is None:
            missing += 1
            continue
        chart = chart_shape.chart
        try:
            cats = [c.label for c in chart.plots[0].categories]
        except Exception:
            cats = []
        if not cats:
            missing += 1
            continue
        v = counts[key]
        _apply_chart(chart, cats, v, answer_key.get(key))
        updated += 1
    prs.save(str(pptx_path))
    return updated, missing


def main():
    if not HAS_PPTX:
        print("pip install python-pptx")
        return 1
    ap = argparse.ArgumentParser(description="Apply poll counts to ITE PowerPoint charts")
    ap.add_argument("--csv", dest="data", type=Path, help="CSV or JSON with vote counts")
    ap.add_argument("--pptx", type=Path, help="Single .pptx file")
    ap.add_argument("--all", action="store_true", help="All topic decks in ITE_PowerPoints/")
    args = ap.parse_args()
    if not args.data or not args.data.exists():
        print("Provide --csv path to poll results (CSV or JSON).")
        print("Generate a template: python export_poll_template.py")
        return 1
    counts = load_counts(args.data)
    answer_key = load_answer_key()
    targets: list[Path] = []
    if args.pptx:
        targets = [args.pptx]
    elif args.all:
        targets = sorted(PPTS_DIR.glob("*.pptx"))
    else:
        print("Use --pptx FILE.pptx or --all")
        return 1
    total_u = total_m = 0
    for p in targets:
        if not p.exists():
            print(f"Skip missing: {p}")
            continue
        u, m = process_presentation(p, counts, answer_key)
        total_u += u
        total_m += m
        print(f"{p.name}: updated {u} poll charts, skipped {m} (no data or no chart)")
    print(f"Done. Total charts updated: {total_u}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
